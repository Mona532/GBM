from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\GBM")
RESULTS = ROOT / "results"
OUTDIR = RESULTS / "ppt"
TEMPLATE = Path(
    r"E:\ILC_review\ppt-master\projects\tls_review_ppt169_20260701\exports\Paper4.pptx"
)


BG = RGBColor(247, 244, 238)
TEXT = RGBColor(28, 28, 28)
MUTED = RGBColor(110, 110, 110)
ACCENT = RGBColor(153, 59, 62)
ACCENT_SOFT = RGBColor(224, 210, 198)
LINE = RGBColor(210, 202, 194)


@dataclass
class FigureSpec:
    title: str
    path: Path
    caption: str


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)  # type: ignore[attr-defined]
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)  # type: ignore[attr-defined]


def pick_blank_layout(prs: Presentation):
    return min(prs.slide_layouts, key=lambda x: len(x.placeholders))


def add_full_background(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.2), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.62), Inches(0.87), Inches(11.0), Inches(0.35))
        p2 = tx2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = MUTED
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.18), Inches(12.1), Pt(1.6)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, text: str = "GBM TLS / 组会汇报 / 2026-07-13") -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.08), Inches(4.0), Inches(0.18))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(8)
    r.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.size


def fit_rect(img_w: int, img_h: int, box_w: float, box_h: float) -> tuple[float, float]:
    ratio = min(box_w / img_w, box_h / img_h)
    return img_w * ratio, img_h * ratio


def add_image(
    slide,
    path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    border: bool = True,
) -> None:
    img_w, img_h = image_size(path)
    target_w, target_h = fit_rect(img_w, img_h, Inches(width), Inches(height))
    x = Inches(left) + (Inches(width) - target_w) / 2
    y = Inches(top) + (Inches(height) - target_h) / 2
    pic = slide.shapes.add_picture(str(path), x, y, width=target_w, height=target_h)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(1.0)


def add_caption(slide, text: str, left: float, top: float, width: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.34))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def add_panel_label(slide, text: str, left: float, top: float, width: float = 1.0) -> None:
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.28)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT_SOFT
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = ACCENT


def build_cover(prs: Presentation, layout) -> None:
    slide = prs.slides.add_slide(layout)
    add_full_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(0.9), Inches(0.28), Inches(5.6)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(10.6), Inches(1.8))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "GBM中TLS异质性解析\n从NMF分型到E4神经免疫耦联"
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT

    sub = slide.shapes.add_textbox(Inches(1.22), Inches(3.1), Inches(7.8), Inches(1.2))
    p2 = sub.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "数据基础：142个空间转录组样本，SpaLinker TLS 区域识别，component-level pseudobulk，17类 cell2loc 注释"
    r2.font.name = "Aptos"
    r2.font.size = Pt(16)
    r2.font.color.rgb = MUTED

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.5), Inches(4.7), Inches(1.35)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_SOFT
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p3 = tf.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "当前主线结论\n1. TLS可分为5类 ecotype\n2. E4呈现明显ILC-受体-神经通路耦联\n3. 空间上可在代表component中直接观察到"
    r3.font.name = "Aptos"
    r3.font.size = Pt(15)
    r3.font.color.rgb = TEXT

    add_footer(slide, "GBM TLS 组会汇报")


def build_outline(prs: Presentation, layout) -> None:
    slide = prs.slides.add_slide(layout)
    add_full_background(slide)
    add_header(slide, "汇报主线", "从 rank-5 ecotype 定义，推进到 E4 的空间与通路证据")
    add_bullets(
        slide,
        [
            "1. 基于 TLS component abundance 的 rank-5 NMF 将 TLS 分成 5 类 ecotype。",
            "2. 5 类 ecotype 在细胞组成、转录组背景和成熟度上都不同。",
            "3. E4 最值得继续追踪，因为其 ILC abundance 与多类神经递质受体呈一致相关。",
            "4. spot-level 相关与 GSEA 进一步支持 E4 内存在神经信号相关转录组背景。",
            "5. 最后回到空间图，确认这种模式确实落在代表性 TLS component 上。",
        ],
        left=0.9,
        top=1.55,
        width=11.0,
        height=4.8,
    )
    add_footer(slide)


def build_single_figure(prs: Presentation, layout, spec: FigureSpec, subtitle: str) -> None:
    slide = prs.slides.add_slide(layout)
    add_full_background(slide)
    add_header(slide, spec.title, subtitle)
    add_image(slide, spec.path, 0.72, 1.45, 11.85, 5.25)
    add_caption(slide, spec.caption, 0.78, 6.78, 11.7)
    add_footer(slide)


def build_two_panel(prs: Presentation, layout, title: str, subtitle: str, left_spec: FigureSpec, right_spec: FigureSpec) -> None:
    slide = prs.slides.add_slide(layout)
    add_full_background(slide)
    add_header(slide, title, subtitle)
    add_panel_label(slide, "左图", 0.76, 1.42, 0.8)
    add_panel_label(slide, "右图", 6.75, 1.42, 0.8)
    add_image(slide, left_spec.path, 0.72, 1.72, 5.7, 4.9)
    add_image(slide, right_spec.path, 6.55, 1.72, 5.7, 4.9)
    add_caption(slide, left_spec.caption, 0.8, 6.72, 5.4)
    add_caption(slide, right_spec.caption, 6.62, 6.72, 5.4)
    add_footer(slide)


def build_five_panel(prs: Presentation, layout, title: str, subtitle: str, specs: list[FigureSpec]) -> None:
    slide = prs.slides.add_slide(layout)
    add_full_background(slide)
    add_header(slide, title, subtitle)
    positions = [
        (0.62, 1.48, 2.35, 2.18),
        (3.02, 1.48, 2.35, 2.18),
        (5.42, 1.48, 2.35, 2.18),
        (7.82, 1.48, 2.35, 2.18),
        (10.22, 1.48, 2.35, 2.18),
    ]
    for i, (spec, pos) in enumerate(zip(specs, positions), start=1):
        left, top, width, height = pos
        add_panel_label(slide, f"P{i}", left + 0.02, top - 0.27, 0.62)
        add_image(slide, spec.path, left, top, width, height)
        add_caption(slide, spec.caption, left + 0.02, top + height + 0.06, width - 0.04)
    add_footer(slide)


def build_summary(prs: Presentation, layout) -> None:
    slide = prs.slides.add_slide(layout)
    add_full_background(slide)
    add_header(slide, "当前可汇报结论", "这版 deck 只基于目前已经落地的主结果，不引入额外推断")
    add_bullets(
        slide,
        [
            "rank-5 NMF 已经把 TLS 组成异质性稳定压缩成 5 类 ecotype，而不只是一个成熟度高低轴。",
            "E4 并不是简单的淋巴富集型，它更像带有 ILC 参与的神经免疫耦联型 TLS 微环境。",
            "受体相关、spot-level 基因相关和 GSEA 三层证据都指向神经递质释放/突触相关程序。",
            "下一步最值得做的是：E4 代表基因/受体的独立队列验证，以及 component 定义的进一步稳健化。",
        ],
        left=0.86,
        top=1.55,
        width=11.0,
        height=4.5,
    )
    note = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(5.65), Inches(11.0), Inches(0.9)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = ACCENT_SOFT
    note.line.fill.background()
    tf = note.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "注：当前 E4 DEG 本身统计学不强，因此本轮汇报避免硬讲 DEG，而改用 ILC abundance 相关的全基因与通路结果。"
    r.font.name = "Aptos"
    r.font.size = Pt(13)
    r.font.color.rgb = TEXT
    add_footer(slide)


def main() -> None:
    OUTDIR.mkdir(parents=True, exist_ok=True)
    template_prs = Presentation(str(TEMPLATE))
    prs = Presentation()
    prs.slide_width = template_prs.slide_width
    prs.slide_height = template_prs.slide_height
    layout = pick_blank_layout(prs)

    build_cover(prs, layout)
    build_outline(prs, layout)

    build_single_figure(
        prs,
        layout,
        FigureSpec(
            "Rank-5 ecotype 定义",
            RESULTS / "fig_tls_rank5_heatmap.jpg",
            "17类 cell2loc abundance 经 row z-score 标准化后可见，E2 近似为全局淋巴富集型，而 E3-E5 提供额外差异轴。",
        ),
        "输入层修正为 17 类 cell2loc 后，TLS abundance basis 不再退化为旧的 13 类结果",
    )

    build_single_figure(
        prs,
        layout,
        FigureSpec(
            "Ecotype 转录组差异热图",
            RESULTS / "fig_tls_rank5_deg_heatmap.jpg",
            "采用 sample-aware 方式控制同一样本内 component 伪重复；这里展示的是 5 类 ecotype 的代表性转录组背景差异。",
        ),
        "重点不是把这些基因硬解释为单细胞 marker，而是把它们视为 component-level transcriptomic context",
    )

    build_single_figure(
        prs,
        layout,
        FigureSpec(
            "成熟度差异热图",
            RESULTS / "fig_tls_rank5_maturity_heatmap.jpg",
            "成熟度信号分布在不同 ecotype 上，并不支持“所有 TLS 共用一个线性成熟度轴”的过度简化。",
        ),
        "maturity 作为独立轴建模，而不是完全由组成丰度替代",
    )

    build_five_panel(
        prs,
        layout,
        "5类 ecotype 的受体表达概览",
        "把不同受体类别拆开看，能更清楚看到每类 ecotype 的受体表达偏好",
        [
            FigureSpec("Glutamate", RESULTS / "fig_tls_rank5_receptor_dotplot_Glutamate.jpg", "Glutamate"),
            FigureSpec("GABA/Gly", RESULTS / "fig_tls_rank5_receptor_dotplot_GABA_Gly.jpg", "GABA/Gly"),
            FigureSpec("Cholinergic", RESULTS / "fig_tls_rank5_receptor_dotplot_Cholinergic.jpg", "Cholinergic"),
            FigureSpec("DA/NE", RESULTS / "fig_tls_rank5_receptor_dotplot_DA_NE.jpg", "DA/NE"),
            FigureSpec("Serotonin", RESULTS / "fig_tls_rank5_receptor_dotplot_Serotonin.jpg", "Serotonin"),
        ],
    )

    build_five_panel(
        prs,
        layout,
        "E4 内 ILC abundance 与受体表达相关性",
        "相关性已改为基于 component 原始 abundance，而不是 component 内相对占比",
        [
            FigureSpec("Glutamate", RESULTS / "fig_tls_rank5_E4_ilc_receptor_heatmap_Glutamate.jpg", "Glutamate"),
            FigureSpec("GABA/Gly", RESULTS / "fig_tls_rank5_E4_ilc_receptor_heatmap_GABA_Gly.jpg", "GABA/Gly"),
            FigureSpec("Cholinergic", RESULTS / "fig_tls_rank5_E4_ilc_receptor_heatmap_Cholinergic.jpg", "Cholinergic"),
            FigureSpec("DA/NE", RESULTS / "fig_tls_rank5_E4_ilc_receptor_heatmap_DA_NE.jpg", "DA/NE"),
            FigureSpec("Serotonin", RESULTS / "fig_tls_rank5_E4_ilc_receptor_heatmap_Serotonin.jpg", "Serotonin"),
        ],
    )

    build_two_panel(
        prs,
        layout,
        "E4 的关键受体和 spot-level 转录组证据",
        "左：HTR1F 与 ILC1/2/3 abundance 的相关散点；右：ILC_total vs gene expression 的摘要相关图",
        FigureSpec(
            "HTR1F 散点图",
            RESULTS / "fig_tls_rank5_E4_HTR1F_scatter.jpg",
            "HTR1F 是当前 E4 里最直观的候选受体之一，用于把 component abundance 相关具体化。",
        ),
        FigureSpec(
            "Figure 2D 风格摘要图",
            RESULTS / "fig_e4_spot_gene_corr_ILC_total.jpg",
            "仅在 E4 component spots 内做全基因相关，避免把非 E4 区域的背景表达混进来。",
        ),
    )

    build_single_figure(
        prs,
        layout,
        FigureSpec(
            "Figure 2E 风格 GSEA",
            RESULTS / "fig_e4_ilc_total_positive_neuro_pathways.jpg",
            "ILC_total 正相关基因的 preranked GSEA 指向 neurotransmitter release、synaptic vesicle exocytosis 等神经相关通路。",
        ),
        "这一步不再讲 DEG，而是讲 E4 中 ILC abundance 增高时，哪些基因程序一起被拉起来",
    )

    build_single_figure(
        prs,
        layout,
        FigureSpec(
            "代表性 E4 空间图",
            RESULTS / "tls_core" / "sample_ecotype_panels" / "E4__AT12-BRA-4-FO-4_1__AT12-BRA-4-FO-4_1_c0_panels.png",
            "完整切片上展示代表性 E4 component 的 TLS score、TLS region 及组成细胞丰度，证明上述模式确实能在空间上落地。",
        ),
        "按样本出图，每个 panel 只高亮目标 component 的代表 spots",
    )

    build_summary(prs, layout)

    output = OUTDIR / "GBM_TLS_group_meeting_20260713.pptx"
    prs.save(str(output))
    print(output)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
